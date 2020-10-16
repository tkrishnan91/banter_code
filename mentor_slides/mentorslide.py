#Purpose of this program is to create a single powerpoint presentation listing all Banter Mentors and their information from mentor survey results.
#Create folder with tsv file from Banter Mentor survey and PTF55F.ttf

from pptx import Presentation

# Open Banter Mentor Survey Results tsv file
alumni = open('mentor.tsv')
readAlumni = alumni.read().split('\n')

# Function definition for creating title slide
# 1) Get title slide layout (builtin layout that has a title and subtitle fields)
# 2) Add a slide to the final presentation with this layout
# 3) Now that the slide is in there, get a handle to the placeholders
# 4) Get a handle to the title placeholder's text text_frame
# 5) call the `para` function to write the appropriate text to title placeholder
# 6) Get a handle to the subtitle placeholder's text text_frame
# 7) call the `para` function to write the appropriate text to the subtitle placeholder


def titleSlide(who):
    title_slidelayout = prs.slide_layouts[0]
    tslide = prs.slides.add_slide(title_slidelayout)
    title = tslide.shapes.placeholders[0]
    subtitle = tslide.shapes.placeholders[1]
    tf = title.text_frame
    para('', "Olin Banter Mentors", tf)
    tf = subtitle.text_frame
    para('', "Alumni Profiles for " + who, tf)

# Create alumni profiles

# Banter Mentor Survey Key in TSV file:
# [0] Timestamp
# [1] What is your preferred email address?
# [2] Name
# [3] Year of Graduation
# [4] Recent Professional Position
# [5] Major at Olin
# [6] I want my mentee to be
# [7] What would you like to get out of Banter?
# [8] What interests and experiences do you have that you could share?
# [9] My interests include...
# [10] At Olin I was involved in...
# [11] These days I am spending time on...
# [12] My coolest project was (or is)...
# [13] My most recent "good find" was ...
# [14] On the weekends I like to...
# [15] If I could travel anywhere I would go to...

# Slide layouts: 0=title, 1=title with textbox, 2=section title
#                3=title with two textboxes, 4=title with three textboxes
#                5=title only, 6=blank


def para(intro, text, tf, fontt='Calibri'):
    text = text.strip()
    if text != "":
        p = tf.add_paragraph()
        p.text = intro + text
#        p.font.name = fontt


# item is a whole row from the spreadsheet
# item is an array of the column values
def profile(item):
    for i in range(len(item)):
        item[i] = item[i].replace('<<<', ',')
    body_slidelayout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(body_slidelayout)

    body_shape = slide.shapes.placeholders[0]  # top
    tf = body_shape.text_frame
    para('', item[2], tf)
    para('', item[5], tf)
    para('', '(' + item[3] + ', ' + item[6] + ')', tf)
    tf.fit_text(font_file='PTF55F.ttf')
    body_shape = slide.shapes.placeholders[2]  # lower left textbox
    tf = body_shape.text_frame

    para("Why Banter: ", item[8], tf)
    para("Areas of interest: ", item[9], tf)
    para("Interests: ", item[10], tf)
    para("Olin Activities: ", item[11], tf)
    para("Current Activities: ", item[12], tf)
    tf.fit_text(font_file='PTF55F.ttf')

    body_shape = slide.shapes.placeholders[4]  # Lower right textbox
    tf = body_shape.text_frame
    para("Recent Good Find: ", item[14], tf)
    para("Travel to: ", item[16], tf)
    para("Coolest Project: ", item[13], tf)
    para("Weekend Activities: ", item[15], tf)
    para("Timezone: ", item[4], tf)
    para('', ' ', tf)  # make sure the last bullet point shows up
    try:
        tf.fit_text(font_file='PTF55F.ttf')
    except:
        print(item)


print ('students')
for item in readAlumni:
    if not item:
        break
    item += '\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t'
    item = item.split('\t')
    # print(item)
    print("Timestamp:" + str(item[0]))
    print("Name:" + str(item[2]) + "\n")
    # print(item)
    if item[7] in ['An Olin Student', "Either (but only one 'mentee)", 'Both (bring it on!)']:
        print("Printing summary")
        print (item[1].replace('<<<', ',')+'---' + item[3]+'---' + item[5].replace('<<<', ','))
# quit()
# Code starts here
rowNumber = 0
# Create a variable called prs that represents Presentation
prs = Presentation()
# Call a titleSlide function to create the title slide
titleSlide('Students and Recent Alumni')

# Read the tsv file; each `item` is a row
for item in readAlumni:
    if not item:
        break
    rowNumber += 1
    if rowNumber == 1:  # skip title row
        continue
    # Add a bunch of tabs to the row (since we are using tsv) and split on tabs
    # This will get us an array of each column entry as an index.
    item += '\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t'
    item = item.split('\t')

    # Pass the array of column values to the profile function
    profile(item)
    # if item[7] in ['An Olin Student', "Either (but only one mentee)", 'Both (bring it on!)']:
    #     profile(item)
    # elif item[7] in ['A younger alum', "Either (but only one mentee)", 'Both (bring it on!)']:
    #     continue
    # else:
    #     print(item)
    if rowNumber > 200:
        break
prs.save('./Mentors_for_Mentees.pptx')
print("should have saved")
