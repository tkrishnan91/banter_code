#Purpose of this program is to create a individual slides listing all Banter Mentees and their information from mentee survey results
#Create folder with tsv file from Banter Mentee survey and PTF55F.ttf

from pptx import Presentation

# Open Banter Mentee Survey Results tsv file
alumni = open('Banter Mentee Signup_ Class of 2017-2020 (Responses) - Form Responses 1.tsv')
readAlumni = alumni.read().split('\n')


# Function definition for creating title slide
# 1) Get title slide layout (builtin layout that has a title and subtitle fields)
# 2) Add a slide to the final presentation with this layout
# 3) Now that the slide is in there, get a handle to the placeholders
# 4) Get a handle to the title placeholder's text text_frame
# 5) call the `para` function to write the appropriate text to title placeholder
# 6) Get a handle to the subtitle placeholder's text text_frame
# 7) call the `para` function to write the appropriate text to the subtitle placeholder

def titleSlide(who):
    title_slidelayout = prs.slide_layouts[0]
    tslide = prs.slides.add_slide(title_slidelayout)
    title = tslide.shapes.placeholders[0]
    subtitle = tslide.shapes.placeholders[1]

    tf = title.text_frame
    para('', "Olin Banter Mentees", tf)
    tf = subtitle.text_frame
    para('', "Mentee Profiles for " + who, tf)


# Create mentee profiles

# Banter Mentee Survey key in TSV file:
# [0]Timestamp
# [1] What is your preferred email address?
# [2] Name
# [3] Year of Graduation
# [4] Probable Major
# [5] What do you want to get out of Banter?
# [6] Anything else you want your match to know?
# [7] My interests include...
# [8] I am involved in...
# [9] If I could travel anywhere I would go to...


# slide layouts: 0=title, 1=title with textbox, 2=section title
#                3=title with two textboxes, 4=title with three textboxes
#                5=title only, 6=blank

def para(intro, text, tf, fontt='Calibri'):
    text = text.strip()
    if text != "":
        p = tf.add_paragraph()
        p.text = intro + text
#        p.font.name = fontt

# item is a whole row from the spreadsheet
# item is an array of the column values

def profile(item, fileCount):
    prs = Presentation()
    for i in range(len(item)):
        item[i] = item[i].replace('<<<',',')
    body_slidelayout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(body_slidelayout)

    body_shape = slide.shapes.placeholders[0] # top
    tf = body_shape.text_frame
    para('', item[2] , tf)
    # para('',item[5], tf)
    # para('','(' + item[3] + ', '+ item[6] +')', tf)
    tf.fit_text(font_file='PTF55F.ttf')


    body_shape = slide.shapes.placeholders[2] # lower left textbox
    tf = body_shape.text_frame

    para("Why Banter: ", item[5], tf)
    para("My interests include: ", item[7], tf)
    para("Current Activities: ", item[8], tf)
    para("Probable Major: ", item[4], tf)
    para("Expected Graduation: ", item[3], tf)
    tf.fit_text(font_file='PTF55F.ttf')

    body_shape = slide.shapes.placeholders[4] # Lower right textbox
    tf = body_shape.text_frame
    para("If I could travel anywhere: ", item[9], tf)
    para("Anything else you want your match to know?: ", item[6], tf)
    # para("Coolest Project: ", item[3], tf)
    # para("Weekend Activities: ", item[15], tf)
    # para("Timezone: ", item[4], tf)
    para('', ' ', tf)  # make sure the last bullet point shows up
    try:
        print("fileCount: " + str(fileCount))
        # save single slide and move on to next
        prs.save(str(fileCount)+"_slide.pptx")
    except:
        print("It threw an exception for " + str(fileCount))


rowNumber=0

# prs = Presentation()
# titleSlide('Banter')

# Read the tsv file; each `item` is a row
for item in readAlumni:
    if not item:
        break
    rowNumber+=1
    if rowNumber == 1: # skip title row
        continue

    # Add a bunch of tabs to the row (since we are using tsv) and split on tabs
    # This will get us an array of each column entry as an index.
    item += '\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t\t'
    item = item.split('\t')


    print("printing item number: " + str(rowNumber))

    # Pass the array of column values to the profile function
    profile(item, rowNumber)

    if rowNumber > 200:
        break
#prs.save('./Mentees_for_Banter.pptx')
print("should have saved")
